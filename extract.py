import os
import re
from pptx import Presentation

pptx_folder = './'
output_txt = 'extract.txt'

# 匹配英文字符和标点符号
english_text_pattern = re.compile(r'[A-Za-z0-9.,!?()";: -]+')

all_english_text = []

# 获取当前目录下所有符合格式的文件
pptx_files = [f for f in os.listdir(pptx_folder) if f.endswith('.pptx') and f[:-5].isdigit()]

# 按照数字顺序排序文件名
pptx_files.sort(key=lambda x: int(x[:-5]))

# 遍历每个 PPT 文件
for pptx_filename in pptx_files:
    pptx_path = os.path.join(pptx_folder, pptx_filename)

    presentation = Presentation(pptx_path)

    # 遍历每一张幻灯片的每个形状
    for slide in presentation.slides:
        for shape in slide.shapes:
            # 确保该形状包含文本框
            if hasattr(shape, 'text'):
                # 获取文本内容
                text = shape.text.strip()

                # 如果文本内容非空，只提取英文文本
                if text:
                    # 用正则表达式提取文本中的英文部分
                    extracted_english_text = ' '.join(english_text_pattern.findall(text))
                    
                    if extracted_english_text:  # 如果有英文文本
                        all_english_text.append(extracted_english_text)

                # 如果文本为空，添加空行
                else:
                    all_english_text.append('')

# 将提取的英文文本写入到 extract.txt 文件
with open(output_txt, 'w', encoding='utf-8') as f:
    for line in all_english_text:
        f.write(line + '\n')

print(f"提取完成，所有英文文本已保存到 {output_txt}")